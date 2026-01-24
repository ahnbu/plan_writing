import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# COLOR PALETTE (Joonggonara / Strategy Theme)
COLOR_PRIMARY = RGBColor(0, 153, 76)      # Joonggonara Green
COLOR_SECONDARY = RGBColor(51, 51, 51)    # Dark Gray (Text)
COLOR_ACCENT = RGBColor(255, 102, 0)      # Orange (Emphasis)
COLOR_BG_LIGHT = RGBColor(245, 245, 245)  # Light Gray for placeholders
COLOR_WHITE = RGBColor(255, 255, 255)

FONT_FAMILY = 'Malgun Gothic'
FONT_FAMILY_BOLD = 'Malgun Gothic' # In case we want a different bold font

def clean_text(text):
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    return text.strip()

def parse_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = []
    slide_blocks = re.split(r'\[슬라이드 \d+\]', content)
    
    for block in slide_blocks:
        if not block.strip():
            continue
            
        slide_data = {
            "title": "",
            "gov_msg": "",
            "body": [],
            "visual": "",
            "type": "Standard" # Default
        }
        
        lines = block.strip().split('\n')
        parsing_body = False
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('- 슬라이드 제목:'):
                parsing_body = False
                title = clean_text(line.replace('- 슬라이드 제목:', '').strip())
                slide_data["title"] = title
                # Simple heuristic for slide type
                if "표지" in title or "Subject" in title:
                    slide_data["type"] = "Title"
                elif "결론" in title or "Conclusion" in title:
                    slide_data["type"] = "Conclusion"
                    
            elif line.startswith('- 거버닝 메시지:'):
                parsing_body = False
                slide_data["gov_msg"] = clean_text(line.replace('- 거버닝 메시지:', '').strip())
            elif line.startswith('- 시각화 계획:'):
                parsing_body = False
                slide_data["visual"] = clean_text(line.replace('- 시각화 계획:', '').strip())
            elif line.startswith('- 본문내용:'):
                parsing_body = True
            elif parsing_body and line.startswith('-'):
                # Handle nested bullets if simple indentation exists
                slide_data["body"].append(clean_text(line.replace('-', '', 1).strip()))
                
        if slide_data["title"]:
            slides.append(slide_data)
            
    return slides

def add_header_footer(slide, slide_num):
    # Add simple footer
    left = Inches(0.5)
    top = SLIDE_HEIGHT - Inches(0.4)
    width = SLIDE_WIDTH - Inches(1.0)
    height = Inches(0.3)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.paragraphs[0]
    p.text = f"Joonggonara Strategy Pivot | {slide_num}"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def create_title_slide(prs, slide_data):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background Graphic (Green Bar on Left)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.8), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()
    
    # Top Logo Area (Placeholder)
    # logo_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(2), Inches(0.5))
    # logo_box.text_frame.text = "JOONGGONARA"
    
    # Title
    top = Inches(2.5)
    left = Inches(1.5)
    width = Inches(10)
    height = Inches(1.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data["title"].replace("표지 : ", "")
    p.font.name = FONT_FAMILY_BOLD
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    
    # Subtitle (Gov Msg)
    top += Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = slide_data["gov_msg"]
    p.font.name = FONT_FAMILY
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_PRIMARY
    
    # Body (Meta info)
    top += Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = txBox.text_frame
    for line in slide_data["body"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)

def create_standard_slide(prs, slide_data, index):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_header_footer(slide, index + 1)
    
    # 1. Header Bar (Top Strip)
    # shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.8), Inches(0.2), Inches(0.5))
    # shape.fill.solid()
    # shape.fill.fore_color.rgb = COLOR_PRIMARY
    # shape.line.fill.background()
    
    # 2. Title
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(12)
    height = Inches(0.8)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.name = FONT_FAMILY_BOLD
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    
    # 3. Governing Message (Kicker Box)
    top = Inches(1.1)
    height = Inches(0.6)
    
    # Background for Msg
    # bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    # bg_shape.fill.solid()
    # bg_shape.fill.fore_color.rgb = RGBColor(235, 245, 235) # Very light green
    # bg_shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.paragraphs[0]
    p.text = slide_data["gov_msg"]
    p.font.name = FONT_FAMILY
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_PRIMARY
    p.font.bold = True
    
    # 4. Content Area (Two Columns)
    # Left: Text Body
    content_top = Inches(2.0)
    col_gap = Inches(0.5)
    left_width = Inches(6.5)
    
    txBox = slide.shapes.add_textbox(left, content_top, left_width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    
    for item in slide_data["body"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(16) # Increased from 14
        p.space_after = Pt(12)
        p.line_spacing = 1.2
        p.level = 0
        
    # Right: Visual Placeholder
    right_left = left + left_width + col_gap
    right_width = SLIDE_WIDTH - right_left - Inches(0.5)
    
    vis_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, content_top, right_width, Inches(4.5)
    )
    vis_shape.fill.solid()
    vis_shape.fill.fore_color.rgb = COLOR_BG_LIGHT
    vis_shape.line.color.rgb = RGBColor(200, 200, 200)
    vis_shape.line.dash_style = 1 # Solid, but maybe we want dashed? (Need import)
    
    tf = vis_shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "[ Visual Plan ]\n\n" + slide_data["visual"]
    p.font.name = FONT_FAMILY
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

def create_conclusion_slide(prs, slide_data, index):
    # Similar to standard but maybe centered?
    # For now, stick to standard but maybe adding an accent
    create_standard_slide(prs, slide_data, index)
    
    # Add an accent bar at the bottom for conclusion
    slide = prs.slides[-1]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(0.2), SLIDE_WIDTH, Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

def create_presentation(slides, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    for i, slide_data in enumerate(slides):
        if i == 0:
            create_title_slide(prs, slide_data)
        elif slide_data["type"] == "Conclusion":
            create_conclusion_slide(prs, slide_data, i)
        else:
            create_standard_slide(prs, slide_data, i)

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    markdown_path = r"d:\Vibe_Coding\plan_writing\output\20260124_Joonggonara\17_slide_plan.md"
    output_pptx = r"d:\Vibe_Coding\plan_writing\output\20260124_Joonggonara\18_Joonggonara_Strategy_Deck_v2.pptx"
    
    slides = parse_markdown(markdown_path)
    create_presentation(slides, output_pptx)
