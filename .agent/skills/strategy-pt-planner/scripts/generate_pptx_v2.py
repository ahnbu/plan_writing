import re
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# COLOR PALETTE (Dynamic via Arguments or Default)
DEFAULT_COLORS = {
    "primary": RGBColor(0, 122, 61),      # Default Green
    "secondary": RGBColor(51, 51, 51),    # Dark Gray
    "accent": RGBColor(238, 30, 48),       # Red
    "orange": RGBColor(248, 150, 31)      # Orange
}

FONT_FAMILY = 'Malgun Gothic'
FONT_FAMILY_BOLD = 'Malgun Gothic'

def clean_text(text):
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    return text.strip()

def parse_markdown(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Markdown file not found: {file_path}")
        
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = []
    # Updated Regex to be more flexible with header formats (e.g., [1], [Slide 1], etc.)
    # It splits primarily by the bracketed number header.
    slide_blocks = re.split(r'\n\s*\[\s*\d+\s*\]', content)
    
    # The first split might be empty or content before the first slide, so check index 0
    # But usually re.split includes the delimiter if captured, here it consumes it.
    # Let's try to capture the block content.
    
    # Alternative: Find all blocks starting with [N]
    # We iterate through the file line by line to build blocks safer.
    
    lines = content.split('\n')
    current_block = []
    
    for line in lines:
        if re.search(r'^\s*\[\s*\d+\s*\]', line):
            if current_block:
                slides.append('\n'.join(current_block))
            current_block = [] # Start new block
        else:
            current_block.append(line)
            
    if current_block:
        slides.append('\n'.join(current_block)) # Add the last block

    parsed_slides = []
    
    for block in slides:
        if not block.strip():
            continue
            
        slide_data = {
            "title": "",
            "gov_msg": "",
            "body": [],
            "visual": "",
            "type": "Standard" 
        }
        
        block_lines = block.strip().split('\n')
        parsing_body = False
        
        for line in block_lines:
            line = line.strip()
            if not line:
                continue
                
            if re.search(r'-\s*\**슬라이드 제목\**:', line):
                parsing_body = False
                title = re.sub(r'-\s*\**슬라이드 제목\**:', '', line).strip()
                slide_data["title"] = clean_text(title)
                if "표지" in slide_data["title"]:
                    slide_data["type"] = "Title"
                elif "결론" in slide_data["title"] or "로드맵" in slide_data["title"]:
                    slide_data["type"] = "Conclusion"
                    
            elif re.search(r'-\s*\**거버닝 메시지\**:', line):
                parsing_body = False
                msg = re.sub(r'-\s*\**거버닝 메시지\**:', '', line).strip()
                slide_data["gov_msg"] = clean_text(msg)
            elif re.search(r'-\s*\**시각화 계획\**:', line):
                parsing_body = False
                vis = re.sub(r'-\s*\**시각화 계획\**:', '', line).strip()
                slide_data["visual"] = clean_text(vis)
            elif re.search(r'-\s*\**본문내용\**:', line):
                parsing_body = True
            elif parsing_body and line.replace('-','',1).strip(): # Check if line has content
                # Handle bullet points
                cleaned_line = line.lstrip(' -').strip()
                if cleaned_line:
                    slide_data["body"].append(clean_text(cleaned_line))
                
        if slide_data["title"]:
            parsed_slides.append(slide_data)
            
    return parsed_slides

def add_header_footer(slide, slide_num, project_name):
    left = Inches(0.5)
    top = SLIDE_HEIGHT - Inches(0.4)
    width = SLIDE_WIDTH - Inches(1.0)
    height = Inches(0.3)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{project_name} | {slide_num}"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.RIGHT

def create_title_slide(prs, slide_data, colors):
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    
    bar_width = Inches(0.3)
    palette = [colors["primary"], colors["accent"], colors["orange"]]
    for i, color in enumerate(palette):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i * 0.3), 0, bar_width, SLIDE_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
    
    top = Inches(2.2)
    left = Inches(1.5)
    width = Inches(11)
    height = Inches(2.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data["title"].split('-', 1)[-1].strip() if '-' in slide_data["title"] else slide_data["title"]
    p.font.name = FONT_FAMILY_BOLD
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors["secondary"]
    
    top += Inches(1.8)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = slide_data["gov_msg"]
    p.font.name = FONT_FAMILY
    p.font.size = Pt(22)
    p.font.color.rgb = colors["primary"]
    
    top += Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = txBox.text_frame
    for line in slide_data["body"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(120, 120, 120)

def create_standard_slide(prs, slide_data, index, colors, project_name):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_header_footer(slide, index + 1, project_name)
    
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12)
    height = Inches(0.8)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.name = FONT_FAMILY_BOLD
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors["secondary"]
    
    top = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = "▶ " + slide_data["gov_msg"]
    p.font.name = FONT_FAMILY
    p.font.size = Pt(16)
    p.font.color.rgb = colors["primary"]
    p.font.bold = True
    
    content_top = Inches(2.2)
    left_width = Inches(7.5)
    
    txBox = slide.shapes.add_textbox(left, content_top, left_width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for item in slide_data["body"]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14)
        p.space_after = Pt(10)
        p.level = 0
        
    vis_left = left + left_width + Inches(0.5)
    vis_width = SLIDE_WIDTH - vis_left - Inches(0.5)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, vis_left, content_top, vis_width, Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "[ Visual Plan ]\n\n" + slide_data["visual"]
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

def main():
    parser = argparse.ArgumentParser(description="Generate Strategy PPTX from Markdown")
    parser.add_argument("--input", required=True, help="Path to input markdown file")
    parser.add_argument("--output", required=True, help="Path to output pptx file")
    parser.add_argument("--project", default="Strategy Report", help="Project name for footer")
    parser.add_argument("--primary", default="0,122,61", help="Primary color (R,G,B)")
    parser.add_argument("--accent", default="238,30,48", help="Accent color (R,G,B)")
    parser.add_argument("--orange", default="248,150,31", help="Orange color (R,G,B)")
    
    args = parser.parse_args()
    
    def parse_rgb(rgb_str):
        r, g, b = map(int, rgb_str.split(','))
        return RGBColor(r, g, b)
        
    colors = {
        "primary": parse_rgb(args.primary),
        "secondary": RGBColor(51, 51, 51),
        "accent": parse_rgb(args.accent),
        "orange": parse_rgb(args.orange)
    }
    
    slides_data = parse_markdown(args.input)
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    for i, slide_data in enumerate(slides_data):
        if slide_data["type"] == "Title":
            create_title_slide(prs, slide_data, colors)
        else:
            create_standard_slide(prs, slide_data, i, colors, args.project)
            
    prs.save(args.output)
    print(f"Presentation saved to {args.output}")

if __name__ == "__main__":
    main()
