import argparse
import sys
import os
from pptx import Presentation

def audit_pptx(pptx_path, plan_path=None):
    if not os.path.exists(pptx_path):
        print(f"❌ Error: File not found - {pptx_path}")
        return False

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"❌ Error: Cannot open PPTX file - {e}")
        return False

    slide_count = len(prs.slides)
    print(f"🔍 Audit Report for: {os.path.basename(pptx_path)}")
    print(f"--------------------------------------------------")
    print(f"Total Slides Found: {slide_count}")

    empty_slides = []
    
    for i, slide in enumerate(prs.slides):
        has_content = False
        title_text = "No Title"
        
        # Check Title
        if slide.shapes.title:
            if slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()
                has_content = True
        
        # Check other shapes for text
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.strip():
                    has_content = True
                    # If title was missing but found text here, good.
        
        print(f"  - Slide {i+1}: {'[OK]' if has_content else '[EMPTY]'} - {title_text[:30]}...")
        
        if not has_content:
            empty_slides.append(i+1)

    print(f"--------------------------------------------------")
    
    # 1. Check Slide Count against Plan (if provided)
    plan_count = 0
    if plan_path and os.path.exists(plan_path):
        with open(plan_path, 'r', encoding='utf-8') as f:
            content = f.read()
            # Count occurrences of "[슬라이드" or similar markers
            import re
            matches = re.findall(r'\[슬라이드 \d+\]', content)
            plan_count = len(matches)
            print(f"Plan Target Slides: {plan_count}")

    # 2. Final Verdict
    is_valid = True
    
    if slide_count <= 1:
        print("❌ FAIL: Too few slides (<= 1). Likely generation error.")
        is_valid = False
    elif slide_count < plan_count * 0.8: # Allow slight deviation
        print(f"❌ FAIL: Slide count ({slide_count}) is significantly less than planned ({plan_count}).")
        is_valid = False
        
    if empty_slides:
        print(f"❌ FAIL: Found empty slides: {empty_slides}")
        is_valid = False

    if is_valid:
        print("✅ PASS: PPTX integrity verified.")
        return True
    else:
        print("🚨 ACTION REQUIRED: Regenerate PPTX.")
        return False

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx_path", required=True)
    parser.add_argument("--plan_path", required=False)
    
    args = parser.parse_args()
    
    success = audit_pptx(args.pptx_path, args.plan_path)
    
    if not success:
        sys.exit(1) # Return error code to trigger task failure
